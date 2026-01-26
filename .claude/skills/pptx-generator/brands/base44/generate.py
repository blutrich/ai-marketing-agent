"""
Base44 PPTX Generator

This module generates on-brand Base44 presentations. It adapts the cookbook
layout patterns to Base44's light theme (gradient background, dark text, orange accents).

Usage:
    from generate import create_presentation, add_title_slide, add_stats_slide, ...

The cookbook layouts are designed for dark themes. This module provides Base44-specific
implementations that follow the same patterns but with the correct light theme colors.
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pathlib import Path


# ============================================================================
# BRAND CONSTANTS - Base44 Light Theme (Updated 2026-01-24)
# ============================================================================

# Colors
BG_TOP = "E8F4F8"           # Light blue (gradient top)
BG_BOTTOM = "FDF5F0"        # Cream/peach (gradient bottom)
BG_SECONDARY = "F0F0F0"     # Alt backgrounds
TEXT = "000000"             # Primary text (true black)
TEXT_SECONDARY = "666666"   # Secondary text (gray)
ACCENT = "FF983B"           # Orange accent (warm orange)
ACCENT_LIGHT = "FFE9DF"     # Light orange backgrounds
ACCENT_MEDIUM = "FFBFA1"    # Secondary accent
ACCENT_DARK = "EA6020"      # Hover/active states
ACCENT_DARKEST = "C94001"   # Deep emphasis
ACCENT_MUTED = "999999"     # Muted gray (for "before" states)
CARD_BG = "FFFFFF"          # White cards
BUTTON_BG = "000000"        # Black buttons
BUTTON_TEXT = "FFFFFF"      # White button text
DIVIDER = "cccccc"          # Light gray dividers

# Fonts (use Arial as fallback for PPTX generation)
# Primary: BC Novatica Cyr / STK Miso (heading), Inter (body)
FONT = "Arial"

# Asset paths
BRAND_DIR = Path(__file__).parent
LOGO_PATH = str(BRAND_DIR / "assets" / "logo.png")

# Slide dimensions
SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5


# ============================================================================
# UTILITY FUNCTIONS
# ============================================================================

def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_gradient_background(slide):
    """Add Base44's signature light gradient background."""
    fill = slide.background.fill
    fill.gradient()
    fill.gradient_angle = 270  # Top to bottom
    fill.gradient_stops[0].color.rgb = hex_to_rgb(BG_TOP)
    fill.gradient_stops[1].color.rgb = hex_to_rgb(BG_BOTTOM)


def add_text_to_shape(shape, text, font_size, font_color, bold=False, alignment=PP_ALIGN.CENTER):
    """Add properly centered text to a shape."""
    tf = shape.text_frame
    tf.word_wrap = False
    tf.anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = hex_to_rgb(font_color)
    p.font.name = FONT
    p.font.bold = bold
    p.alignment = alignment


def add_textbox(slide, text, x, y, width, height, font_size, color=TEXT,
                bold=False, alignment=PP_ALIGN.LEFT, word_wrap=True):
    """Add a textbox with consistent styling."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = hex_to_rgb(color)
    p.font.name = FONT
    p.font.bold = bold
    p.alignment = alignment
    return tb


def add_accent_bar(slide, x, y, width, height, color=ACCENT):
    """Add an accent bar."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(color)
    bar.line.fill.background()
    return bar


# ============================================================================
# PRESENTATION FACTORY
# ============================================================================

def create_presentation():
    """Create a new Base44-branded presentation."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)
    return prs


# ============================================================================
# SLIDE TYPES
# ============================================================================

def add_title_slide(prs, headline, subtitle=None):
    """
    Add a title slide.

    Args:
        prs: Presentation object
        headline: Main headline text
        subtitle: Optional subtitle text
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide)

    # Logo
    slide.shapes.add_picture(LOGO_PATH, Inches(0.6), Inches(0.5), width=Inches(2))

    # Accent bar
    add_accent_bar(slide, 0.6, 2.6, 1.2, 0.06)

    # Headline
    add_textbox(slide, headline, 0.6, 2.9, 10, 1, 56, TEXT, bold=True)

    # Subtitle
    if subtitle:
        add_textbox(slide, subtitle, 0.6, 4.1, 9, 0.6, 22, TEXT_SECONDARY)

    return slide


def add_stats_slide(prs, title, stats):
    """
    Add a stats/metrics slide.

    Args:
        prs: Presentation object
        title: Slide title
        stats: List of (value, label) tuples, 2-4 items
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide)

    # Title
    add_textbox(slide, title, 0.5, 0.5, 12, 0.8, 40, TEXT, bold=True)

    # Stats layout
    num_stats = len(stats)
    total_width = 12.0
    stat_width = total_width / num_stats
    start_x = 0.66

    for i, (value, label) in enumerate(stats):
        x = start_x + (i * stat_width)

        # Value (big number)
        add_textbox(slide, value, x, 2.3, stat_width - 0.3, 1.5, 64, ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

        # Underline accent
        add_accent_bar(slide, x + 0.8, 4.0, stat_width - 1.9, 0.06)

        # Label
        add_textbox(slide, label, x, 4.3, stat_width - 0.3, 0.8, 18, TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

    # Dividers between stats
    for i in range(1, num_stats):
        x = start_x + (i * stat_width) - 0.15
        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.5), Inches(0.02), Inches(2.2))
        div.fill.solid()
        div.fill.fore_color.rgb = hex_to_rgb(DIVIDER)
        div.line.fill.background()

    return slide


def add_two_column_slide(prs, title, left_header, left_bullets, right_header, right_bullets,
                         left_color=ACCENT_MUTED, right_color=ACCENT):
    """
    Add a two-column comparison slide.

    Args:
        prs: Presentation object
        title: Slide title
        left_header: Left column header
        left_bullets: List of bullet points for left column
        right_header: Right column header
        right_bullets: List of bullet points for right column
        left_color: Color for left column indicators (default: muted gray)
        right_color: Color for right column indicators (default: accent orange)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide)

    # Title
    add_textbox(slide, title, 0.5, 0.5, 12, 0.8, 40, TEXT, bold=True)

    # Center divider
    add_accent_bar(slide, 6.6, 1.5, 0.06, 5.0, ACCENT)

    # LEFT SIDE
    # Header indicator
    add_accent_bar(slide, 0.5, 1.7, 0.12, 0.45, left_color)

    # Header text
    add_textbox(slide, left_header, 0.8, 1.65, 5.5, 0.5, 24, left_color, bold=True)

    # Bullets
    for i, bullet in enumerate(left_bullets):
        y = 2.5 + (i * 1.0)
        # Bullet marker
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y + 0.1), Inches(0.1), Inches(0.1))
        marker.fill.solid()
        marker.fill.fore_color.rgb = hex_to_rgb(left_color)
        marker.line.fill.background()
        # Bullet text
        add_textbox(slide, bullet, 0.8, y, 5.5, 0.8, 18, TEXT)

    # RIGHT SIDE
    # Header indicator
    add_accent_bar(slide, 6.9, 1.7, 0.12, 0.45, right_color)

    # Header text
    add_textbox(slide, right_header, 7.2, 1.65, 5.5, 0.5, 24, right_color, bold=True)

    # Bullets
    for i, bullet in enumerate(right_bullets):
        y = 2.5 + (i * 1.0)
        # Bullet marker
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.9), Inches(y + 0.1), Inches(0.1), Inches(0.1))
        marker.fill.solid()
        marker.fill.fore_color.rgb = hex_to_rgb(right_color)
        marker.line.fill.background()
        # Bullet text
        add_textbox(slide, bullet, 7.2, y, 5.5, 0.8, 18, TEXT)

    return slide


def add_quote_slide(prs, quote, attribution):
    """
    Add a quote slide.

    Args:
        prs: Presentation object
        quote: The quote text (without quotation marks)
        attribution: Who said it (e.g., "— Name, Title")
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide)

    # Accent bar on left
    add_accent_bar(slide, 1.0, 2.5, 0.12, 2.5)

    # Quote text (with inline quotation marks)
    add_textbox(slide, f'"{quote}"', 1.5, 2.5, 10.5, 2.5, 28, TEXT)

    # Attribution
    add_textbox(slide, attribution, 1.5, 5.3, 10, 0.5, 18, TEXT_SECONDARY)

    return slide


def add_cta_slide(prs, headline, subtext=None, button_text="base44.com"):
    """
    Add a CTA/closing slide.

    Args:
        prs: Presentation object
        headline: Main headline
        subtext: Optional supporting text
        button_text: Button label
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide)

    # Centered logo
    logo_w = 2.2
    slide.shapes.add_picture(LOGO_PATH, Inches((SLIDE_WIDTH - logo_w) / 2), Inches(1.5), width=Inches(logo_w))

    # Headline
    add_textbox(slide, headline, 0.5, 3.2, 12.333, 0.9, 48, TEXT, bold=True, alignment=PP_ALIGN.CENTER)

    # Subtext
    if subtext:
        add_textbox(slide, subtext, 0.5, 4.2, 12.333, 0.5, 20, TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

    # Button
    btn_w, btn_h = 3.0, 0.6
    btn_x = (SLIDE_WIDTH - btn_w) / 2
    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(btn_x), Inches(5.2), Inches(btn_w), Inches(btn_h))
    btn.fill.solid()
    btn.fill.fore_color.rgb = hex_to_rgb(BUTTON_BG)
    btn.line.fill.background()
    add_text_to_shape(btn, button_text, 18, BUTTON_TEXT, bold=True)

    return slide


def add_content_slide(prs, title, bullets):
    """
    Add a content slide with bullet points.

    Args:
        prs: Presentation object
        title: Slide title
        bullets: List of bullet points (3-6 items)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide)

    # Title
    add_textbox(slide, title, 0.5, 0.5, 12, 0.8, 40, TEXT, bold=True)

    # Accent bar under title
    add_accent_bar(slide, 0.5, 1.4, 2.0, 0.06)

    # Bullets
    for i, bullet in enumerate(bullets):
        y = 2.0 + (i * 0.9)
        # Bullet marker
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y + 0.12), Inches(0.12), Inches(0.12))
        marker.fill.solid()
        marker.fill.fore_color.rgb = hex_to_rgb(ACCENT)
        marker.line.fill.background()
        # Bullet text
        add_textbox(slide, bullet, 0.85, y, 11.5, 0.7, 20, TEXT)

    return slide


def add_cards_slide(prs, title, cards):
    """
    Add a slide with 3 feature cards.

    Args:
        prs: Presentation object
        title: Slide title
        cards: List of (card_title, card_description) tuples (exactly 3)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide)

    # Title
    add_textbox(slide, title, 0.5, 0.5, 12, 0.8, 40, TEXT, bold=True)

    # Cards layout
    card_width = 3.8
    card_height = 3.5
    gap = 0.35
    total_width = (card_width * 3) + (gap * 2)
    start_x = (SLIDE_WIDTH - total_width) / 2
    y = 1.8

    for i, (card_title, card_desc) in enumerate(cards):
        x = start_x + i * (card_width + gap)

        # Card shadow
        shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x + 0.05), Inches(y + 0.08), Inches(card_width), Inches(card_height))
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = hex_to_rgb(DIVIDER)
        shadow.line.fill.background()

        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(card_width), Inches(card_height))
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(CARD_BG)
        card.line.fill.background()

        # Accent bar
        add_accent_bar(slide, x + 0.3, y + 0.3, 0.6, 0.06)

        # Card title
        add_textbox(slide, card_title, x + 0.3, y + 0.5, card_width - 0.6, 0.6, 20, TEXT, bold=True)

        # Card description
        add_textbox(slide, card_desc, x + 0.3, y + 1.2, card_width - 0.6, 2.0, 16, TEXT_SECONDARY)

    return slide


# ============================================================================
# CONVENIENCE FUNCTION
# ============================================================================

def save_presentation(prs, filename, output_dir=None):
    """
    Save the presentation.

    Args:
        prs: Presentation object
        filename: Filename (without extension)
        output_dir: Output directory (default: Base44 Research output folder)

    Returns:
        Path to saved file
    """
    if output_dir is None:
        output_dir = Path("/Users/oferbl/Desktop/Dev/Base44 Research/output")
    else:
        output_dir = Path(output_dir)

    output_dir.mkdir(parents=True, exist_ok=True)

    output_path = output_dir / f"{filename}.pptx"
    prs.save(output_path)
    return output_path
