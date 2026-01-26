"""
Base44 PPTX Generation Helpers
Reusable functions for creating properly aligned, on-brand slides.
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pathlib import Path

# ============ BRAND CONSTANTS (Updated 2026-01-24) ============

# Colors (Light Theme - New Palette)
BRAND_BG_TOP = "E8F4F8"         # Light blue
BRAND_BG_BOTTOM = "FDF5F0"      # Cream/peach
BRAND_BG_SECONDARY = "F0F0F0"   # Alt backgrounds
BRAND_TEXT = "000000"           # Primary text (true black)
BRAND_TEXT_SECONDARY = "666666"  # Secondary text
BRAND_ACCENT = "FF983B"         # Orange accent (warm orange)
BRAND_ACCENT_LIGHT = "FFE9DF"   # Light orange backgrounds
BRAND_ACCENT_MEDIUM = "FFBFA1"  # Secondary accent
BRAND_ACCENT_DARK = "EA6020"    # Hover/active states
BRAND_ACCENT_DARKEST = "C94001" # Deep emphasis
BRAND_CARD_BG = "FFFFFF"        # White cards
BRAND_BUTTON_BG = "000000"      # Black buttons
BRAND_BUTTON_TEXT = "FFFFFF"    # White button text

# Fonts (use Arial as fallback for PPTX generation)
# Primary: BC Novatica Cyr / STK Miso (heading), Inter (body)
BRAND_HEADING_FONT = "Arial"
BRAND_BODY_FONT = "Arial"

# Paths
BRAND_DIR = Path(__file__).parent
LOGO_PATH = BRAND_DIR / "assets" / "logo.png"

# Slide dimensions
SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5


# ============ HELPER FUNCTIONS ============

def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def create_presentation() -> Presentation:
    """Create a new presentation with Base44 dimensions."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)
    return prs


def add_gradient_background(slide, color_top=BRAND_BG_TOP, color_bottom=BRAND_BG_BOTTOM):
    """Add Base44 gradient background to slide."""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 270  # Top to bottom
    fill.gradient_stops[0].color.rgb = hex_to_rgb(color_top)
    fill.gradient_stops[1].color.rgb = hex_to_rgb(color_bottom)


def add_text_to_shape(shape, text, font_size, font_color=BRAND_TEXT,
                      font_name=BRAND_BODY_FONT, bold=False,
                      alignment=PP_ALIGN.CENTER):
    """
    Add properly centered text to a shape.
    Text is both horizontally and vertically centered.
    """
    tf = shape.text_frame
    tf.word_wrap = False
    tf.auto_size = None
    tf.anchor = MSO_ANCHOR.MIDDLE  # Vertical center

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = hex_to_rgb(font_color)
    p.font.name = font_name
    p.font.bold = bold
    p.alignment = alignment


def add_logo(slide, x=0.6, y=0.5, width=2):
    """Add Base44 logo to slide."""
    return slide.shapes.add_picture(str(LOGO_PATH), Inches(x), Inches(y), width=Inches(width))


def add_logo_centered(slide, y=1.5, width=2.2):
    """Add Base44 logo centered horizontally."""
    x = (SLIDE_WIDTH - width) / 2
    return slide.shapes.add_picture(str(LOGO_PATH), Inches(x), Inches(y), width=Inches(width))


def add_headline(slide, text, x=0.6, y=2.8, width=10, font_size=60,
                 alignment=PP_ALIGN.LEFT):
    """Add a headline text box."""
    title_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.font.name = BRAND_HEADING_FONT
    p.alignment = alignment
    return title_box


def add_subtitle(slide, text, x=0.6, y=4.2, width=9, font_size=24,
                 alignment=PP_ALIGN.LEFT):
    """Add a subtitle text box."""
    sub_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.8))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
    p.font.name = BRAND_BODY_FONT
    p.alignment = alignment
    return sub_box


def add_accent_bar(slide, x=0.6, y=2.5, width=1.5, height=0.08):
    """Add orange accent bar."""
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y),
        Inches(width), Inches(height)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    accent.line.fill.background()
    return accent


def add_button(slide, text, x=None, y=4.8, width=3.2, height=0.65, centered=True):
    """
    Add a black CTA button with properly centered text.
    If centered=True and x is None, button is horizontally centered.
    """
    if centered and x is None:
        x = (SLIDE_WIDTH - width) / 2

    button = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y),
        Inches(width), Inches(height)
    )
    button.fill.solid()
    button.fill.fore_color.rgb = hex_to_rgb(BRAND_BUTTON_BG)
    button.line.fill.background()

    # Add text directly to shape (properly centered)
    add_text_to_shape(
        button,
        text,
        font_size=20,
        font_color=BRAND_BUTTON_TEXT,
        font_name=BRAND_BODY_FONT,
        bold=True
    )
    return button


def add_card(slide, title, description, x, y, width=3.8, height=3.2):
    """
    Add a white card with shadow, accent bar, title, and description.
    All elements are properly spaced and aligned.
    """
    # Card shadow
    shadow = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x + 0.05), Inches(y + 0.08),
        Inches(width), Inches(height)
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(200, 200, 200)
    shadow.line.fill.background()

    # White card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = hex_to_rgb(BRAND_CARD_BG)
    card.line.fill.background()

    # Orange accent bar
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x + 0.3), Inches(y + 0.3),
        Inches(0.6), Inches(0.06)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    accent_bar.line.fill.background()

    # Card title
    title_tb = slide.shapes.add_textbox(
        Inches(x + 0.3), Inches(y + 0.5),
        Inches(width - 0.6), Inches(0.6)
    )
    tf = title_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.font.name = BRAND_HEADING_FONT

    # Card description
    desc_tb = slide.shapes.add_textbox(
        Inches(x + 0.3), Inches(y + 1.2),
        Inches(width - 0.6), Inches(1.8)
    )
    tf = desc_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(16)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
    p.font.name = BRAND_BODY_FONT

    return card


def add_cards_row(slide, cards_data, y=1.7, card_width=3.8, card_height=3.2, gap=0.35):
    """
    Add a row of cards. cards_data is list of (title, description) tuples.
    Supports 2-4 cards, automatically centered.
    """
    num_cards = len(cards_data)
    total_width = num_cards * card_width + (num_cards - 1) * gap
    start_x = (SLIDE_WIDTH - total_width) / 2

    for i, (title, desc) in enumerate(cards_data):
        x = start_x + i * (card_width + gap)
        add_card(slide, title, desc, x, y, card_width, card_height)


# ============ SLIDE TEMPLATES ============

def create_title_slide(prs, headline, subtitle=None):
    """Create a title slide with logo, headline, and optional subtitle."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide)
    add_logo(slide)
    add_accent_bar(slide)
    add_headline(slide, headline)
    if subtitle:
        add_subtitle(slide, subtitle)
    return slide


def create_cards_slide(prs, title, cards_data):
    """Create a slide with title and row of feature cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(11.733), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.font.name = BRAND_HEADING_FONT

    # Cards
    add_cards_row(slide, cards_data)
    return slide


def create_cta_slide(prs, headline, button_text="base44.com"):
    """Create a CTA slide with centered logo, headline, and button."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide)
    add_logo_centered(slide, y=1.8)
    add_headline(slide, headline, x=0.6, y=3.5, width=12.133, alignment=PP_ALIGN.CENTER)
    add_button(slide, button_text)
    return slide
