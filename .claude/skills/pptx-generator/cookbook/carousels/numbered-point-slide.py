#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "carousel-numbered-point-slide"
# format = "carousel"
# dimensions = "square"
# purpose = "Numbered list item with big number - for listicle carousels"
# best_for = [
#     "Listicle content (5 tips, 7 mistakes, etc.)",
#     "Step-by-step processes",
#     "Ranked items",
#     "Sequential content where order matters",
# ]
# avoid_when = [
#     "Standalone points without sequence - use single-point-slide",
#     "First slide - use hook-slide with the list premise",
#     "Last slide - use cta-slide",
# ]
# point_text_max_chars = 60
# supporting_max_chars = 100
# instructions = [
#     "NUMBER should be the position (1, 2, 3, etc.)",
#     "POINT_TEXT is the main content - max 60 chars",
#     "SUPPORTING_TEXT is optional explanation - max 100 chars",
#     "Keep consistent numbering style across carousel",
#     "Logo appears at bottom for brand recognition",
# ]
# ///
"""
LAYOUT: Carousel Numbered Point Slide (Base44 Light Theme)
PURPOSE: Single numbered item from a list - big number with point
FORMAT: Square (1:1) for LinkedIn carousels

CUSTOMIZE:
- BRAND_* colors from brand.json
- NUMBER: The list position (1, 2, 3, etc.)
- POINT_TEXT: The main point (max 60 chars)
- SUPPORTING_TEXT: Optional explanation (max 100 chars)
- LOGO_PATH: Path to brand logo
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_logo(slide, logo_path: str, x: float, y: float, width: float = 1.2):
    """Add logo to slide if file exists."""
    logo_file = Path(logo_path)
    if logo_file.exists():
        slide.shapes.add_picture(
            str(logo_file),
            Inches(x), Inches(y),
            width=Inches(width)
        )
        return True
    return False


def main() -> None:
    # === BASE44 BRAND COLORS (Light Theme - Updated 2026-01-24) ===
    BRAND_BG = "E8F4F8"           # Light blue
    BRAND_TEXT = "000000"         # True black text
    BRAND_TEXT_SECONDARY = "666666"
    BRAND_ACCENT = "FF983B"       # Warm orange
    BRAND_CARD_BG = "FFFFFF"      # White
    BRAND_HEADING_FONT = "Arial"  # Fallback for BC Novatica Cyr
    BRAND_BODY_FONT = "Arial"     # Fallback for Inter

    # === CONTENT ===
    NUMBER = "1"  # Position in list
    POINT_TEXT = "REPLACE"  # Max 60 chars
    SUPPORTING_TEXT = ""  # Optional, max 100 chars
    LOGO_PATH = "brands/base44/assets/logo.png"

    # === SLIDE SETUP (Square 1:1) ===
    prs = Presentation()
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background - light blue
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # White content card - takes up most of slide
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.35), Inches(0.35),
        Inches(6.8), Inches(6.8)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = hex_to_rgb(BRAND_CARD_BG)
    card.line.fill.background()
    card.adjustments[0] = 0.03

    # Orange accent bar at top
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.35), Inches(0.35),
        Inches(6.8), Inches(0.1)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    accent_bar.line.fill.background()

    # Orange number circle
    num_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(0.8), Inches(1.2),
        Inches(1.5), Inches(1.5)
    )
    num_circle.fill.solid()
    num_circle.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    num_circle.line.fill.background()

    # Number text inside circle - centered
    num_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.2),
        Inches(1.5), Inches(1.5)
    )
    tf = num_box.text_frame
    tf.anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].text = NUMBER
    tf.paragraphs[0].font.name = BRAND_HEADING_FONT
    tf.paragraphs[0].font.size = Pt(64)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = hex_to_rgb("FFFFFF")
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Point text - centered in remaining space
    point_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(3.0),
        Inches(6.1), Inches(2.0)
    )
    tf = point_box.text_frame
    tf.word_wrap = True
    tf.anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = POINT_TEXT
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.CENTER

    # Supporting text (if provided)
    if SUPPORTING_TEXT:
        support_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(5.0),
            Inches(6.1), Inches(1.0)
        )
        tf = support_box.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.text = SUPPORTING_TEXT
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
        p.alignment = PP_ALIGN.CENTER

    # NOTE: Logo REMOVED from content slide - logo only on CTA slide
    # This fixes Mac Quick Look rendering bug where text stops rendering
    # when same image appears on multiple slides

    # Save
    output = Path("carousel-numbered-point-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
