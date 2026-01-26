#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "carousel-single-point-slide"
# format = "carousel"
# dimensions = "square"
# purpose = "One key point with supporting text - the workhorse slide"
# best_for = [
#     "Main content slides in carousel body",
#     "Single concept that needs explanation",
#     "Key takeaways or insights",
#     "Standalone valuable points",
# ]
# avoid_when = [
#     "First slide - use hook-slide instead",
#     "Last slide - use cta-slide instead",
#     "Numbered lists - use numbered-point-slide instead",
#     "Quotes - use quote-slide instead",
# ]
# headline_max_chars = 40
# body_max_chars = 150
# instructions = [
#     "Headline should be the KEY POINT in 40 chars or less",
#     "Body text expands on the headline - max 150 chars",
#     "One idea per slide - if you have multiple, split them",
#     "Make headline readable at thumbnail size",
#     "Logo appears at bottom for brand recognition",
# ]
# ///
"""
LAYOUT: Carousel Single Point Slide (Base44 Light Theme)
PURPOSE: One key point with supporting explanation
FORMAT: Square (1:1) for LinkedIn carousels

CUSTOMIZE:
- BRAND_* colors from brand.json
- HEADLINE: The main point (max 40 chars)
- BODY_TEXT: Supporting explanation (max 150 chars)
- LOGO_PATH: Path to brand logo
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_logo(slide, logo_path: str, x: float, y: float, width: float = 1.2):
    """Add logo to slide if file exists."""
    logo_file = Path(logo_path)
    if logo_file.exists():
        slide.shapes.add_picture(
            str(logo_file),
            Inches(x), Inches(y),
            width=Inches(width)
        )
        return True
    return False


def main() -> None:
    # === BASE44 BRAND COLORS (Light Theme - Updated 2026-01-24) ===
    BRAND_BG = "E8F4F8"           # Light blue
    BRAND_TEXT = "000000"         # True black text
    BRAND_TEXT_SECONDARY = "666666"
    BRAND_ACCENT = "FF983B"       # Warm orange
    BRAND_CARD_BG = "FFFFFF"      # White
    BRAND_HEADING_FONT = "Arial"  # Fallback for BC Novatica Cyr
    BRAND_BODY_FONT = "Arial"     # Fallback for Inter

    # === CONTENT ===
    HEADLINE = "REPLACE"  # Max 40 chars
    BODY_TEXT = "REPLACE"  # Max 150 chars
    LOGO_PATH = "brands/base44/assets/logo.png"

    # === SLIDE SETUP (Square 1:1) ===
    prs = Presentation()
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background - light blue
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # White content card - takes up most of slide
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.35), Inches(0.35),
        Inches(6.8), Inches(6.8)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = hex_to_rgb(BRAND_CARD_BG)
    card.line.fill.background()
    card.adjustments[0] = 0.03

    # Orange accent bar at top of card
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.35), Inches(0.35),
        Inches(6.8), Inches(0.1)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    accent_bar.line.fill.background()

    # Content container - centered vertically in card
    # Headline
    headline_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5),
        Inches(6.1), Inches(2.5)
    )
    tf = headline_box.text_frame
    tf.word_wrap = True
    tf.anchor = MSO_ANCHOR.BOTTOM  # Align to bottom of box
    p = tf.paragraphs[0]
    p.text = HEADLINE
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.CENTER

    # Orange accent line (centered)
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.0), Inches(4.1),
        Inches(1.5), Inches(0.06)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    accent_line.line.fill.background()

    # Body text
    body_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(4.3),
        Inches(6.1), Inches(1.8)
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.anchor = MSO_ANCHOR.TOP  # Align to top of box
    p = tf.paragraphs[0]
    p.text = BODY_TEXT
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(22)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
    p.alignment = PP_ALIGN.CENTER

    # NOTE: Logo REMOVED from content slide - logo only on CTA slide
    # This fixes Mac Quick Look rendering bug where text stops rendering
    # when same image appears on multiple slides

    # Save
    output = Path("carousel-single-point-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
