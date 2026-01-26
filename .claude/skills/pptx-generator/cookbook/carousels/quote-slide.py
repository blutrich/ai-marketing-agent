#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "carousel-quote-slide"
# format = "carousel"
# dimensions = "square"
# purpose = "Quote with attribution - for social proof or memorable statements"
# best_for = [
#     "Social proof from customers or experts",
#     "Memorable statements worth highlighting",
#     "Your own quotable insights",
#     "Breaking up content with a different format",
# ]
# avoid_when = [
#     "Long quotes over 120 characters - paraphrase instead",
#     "Quotes that need context - add context in previous slide",
#     "First slide - use hook-slide instead",
# ]
# quote_max_chars = 120
# attribution_max_chars = 40
# instructions = [
#     "Quote should be IMPACTFUL and max 120 characters",
#     "Attribution is optional but recommended for credibility",
#     "Use quotation marks in the design, not the text",
#     "If it's your own quote, attribution can be your handle",
#     "Logo appears at bottom for brand recognition",
# ]
# ///
"""
LAYOUT: Carousel Quote Slide (Base44 Light Theme)
PURPOSE: Featured quote with attribution
FORMAT: Square (1:1) for LinkedIn carousels

CUSTOMIZE:
- BRAND_* colors from brand.json
- QUOTE_TEXT: The quote (max 120 chars, no quotation marks)
- ATTRIBUTION: Who said it (max 40 chars, optional)
- LOGO_PATH: Path to brand logo
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_logo(slide, logo_path: str, x: float, y: float, width: float = 1.2):
    """Add logo to slide if file exists."""
    logo_file = Path(logo_path)
    if logo_file.exists():
        slide.shapes.add_picture(
            str(logo_file),
            Inches(x), Inches(y),
            width=Inches(width)
        )
        return True
    return False


def main() -> None:
    # === BASE44 BRAND COLORS (Light Theme - Updated 2026-01-24) ===
    BRAND_BG = "E8F4F8"           # Light blue
    BRAND_TEXT = "000000"         # True black text
    BRAND_ACCENT = "FF983B"       # Warm orange
    BRAND_CARD_BG = "FFFFFF"      # White
    BRAND_HEADING_FONT = "Arial"  # Fallback for BC Novatica Cyr
    BRAND_BODY_FONT = "Arial"     # Fallback for Inter

    # === CONTENT ===
    QUOTE_TEXT = "REPLACE"  # Max 120 chars, no quotation marks
    ATTRIBUTION = ""  # Optional, max 40 chars
    LOGO_PATH = "brands/base44/assets/logo.png"

    # === SLIDE SETUP (Square 1:1) ===
    prs = Presentation()
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background - light blue
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # White content card - takes up most of slide
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.35), Inches(0.35),
        Inches(6.8), Inches(6.8)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = hex_to_rgb(BRAND_CARD_BG)
    card.line.fill.background()
    card.adjustments[0] = 0.03

    # Orange accent bar at top of card
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.35), Inches(0.35),
        Inches(6.8), Inches(0.1)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    accent_bar.line.fill.background()

    # Large opening quote mark (orange)
    open_quote = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.8),
        Inches(1.5), Inches(1.5)
    )
    p = open_quote.text_frame.paragraphs[0]
    p.text = """
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

    # Quote text - centered both ways
    quote_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(2.2),
        Inches(6.1), Inches(3.0)
    )
    tf = quote_box.text_frame
    tf.word_wrap = True
    tf.anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = QUOTE_TEXT
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(26)
    p.font.italic = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.CENTER

    # Attribution
    if ATTRIBUTION:
        attr_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(5.3),
            Inches(6.1), Inches(0.5)
        )
        tf = attr_box.text_frame
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = f"— {ATTRIBUTION}"
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
        p.alignment = PP_ALIGN.CENTER

    # NOTE: Logo REMOVED from content slide - logo only on CTA slide
    # This fixes Mac Quick Look rendering bug where text stops rendering
    # when same image appears on multiple slides

    # Save
    output = Path("carousel-quote-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
