#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "carousel-cta-slide"
# format = "carousel"
# dimensions = "square"
# purpose = "Closing slide with call to action - the conversion slide"
# best_for = [
#     "Last slide of any carousel",
#     "Follow/subscribe prompts",
#     "Link to resource or content",
#     "Summary with next step",
# ]
# avoid_when = [
#     "First or middle slides - save CTA for the end",
#     "Multiple CTAs - pick ONE action",
#     "No clear action - every carousel needs a CTA",
# ]
# headline_max_chars = 30
# cta_max_chars = 25
# instructions = [
#     "HEADLINE should summarize value or prompt action - max 30 chars",
#     "CTA_TEXT is the specific action - max 25 chars (e.g., 'Try Base44 Free')",
#     "Logo appears prominently for brand recognition",
#     "Keep it simple - one clear action beats multiple options",
# ]
# ///
"""
LAYOUT: Carousel CTA Slide (Base44 Light Theme)
PURPOSE: Closing slide with call to action
FORMAT: Square (1:1) for LinkedIn carousels

CUSTOMIZE:
- BRAND_* colors from brand.json
- HEADLINE: Summary or prompt (max 30 chars)
- CTA_TEXT: The action (max 25 chars)
- LOGO_PATH: Path to brand logo
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_logo(slide, logo_path: str, x: float, y: float, width: float = 1.5):
    """Add logo to slide if file exists."""
    logo_file = Path(logo_path)
    if logo_file.exists():
        slide.shapes.add_picture(
            str(logo_file),
            Inches(x), Inches(y),
            width=Inches(width)
        )
        return True
    return False


def main() -> None:
    # === BASE44 BRAND COLORS (Light Theme - Updated 2026-01-24) ===
    BRAND_BG = "E8F4F8"           # Light blue
    BRAND_TEXT = "000000"         # True black text
    BRAND_TEXT_SECONDARY = "666666"
    BRAND_ACCENT = "FF983B"       # Warm orange
    BRAND_CARD_BG = "FFFFFF"      # White
    BRAND_BUTTON_BG = "000000"    # Black buttons
    BRAND_HEADING_FONT = "Arial"  # Fallback for BC Novatica Cyr
    BRAND_BODY_FONT = "Arial"     # Fallback for Inter

    # === CONTENT ===
    HEADLINE = "REPLACE"  # Max 30 chars
    CTA_TEXT = "REPLACE"  # Max 25 chars
    LOGO_PATH = "brands/base44/assets/logo.png"

    # === SLIDE SETUP (Square 1:1) ===
    prs = Presentation()
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background - light blue
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # White content card - takes up most of slide
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.35), Inches(0.35),
        Inches(6.8), Inches(6.8)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = hex_to_rgb(BRAND_CARD_BG)
    card.line.fill.background()
    card.adjustments[0] = 0.03

    # Orange accent bar at top
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.35), Inches(0.35),
        Inches(6.8), Inches(0.1)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    accent_bar.line.fill.background()

    # Logo - centered and prominent
    add_logo(slide, LOGO_PATH, x=3.0, y=1.3, width=1.5)

    # Headline - centered
    headline_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.0),
        Inches(6.5), Inches(1.2)
    )
    tf = headline_box.text_frame
    tf.word_wrap = True
    tf.anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = HEADLINE
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.CENTER

    # CTA Button (black rounded rectangle)
    cta_button = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.75), Inches(4.5),
        Inches(4.0), Inches(0.9)
    )
    cta_button.fill.solid()
    cta_button.fill.fore_color.rgb = hex_to_rgb(BRAND_BUTTON_BG)
    cta_button.line.fill.background()
    cta_button.adjustments[0] = 0.5  # Fully rounded

    # CTA text - centered in button
    cta_tf = cta_button.text_frame
    cta_tf.anchor = MSO_ANCHOR.MIDDLE
    cta_tf.paragraphs[0].text = CTA_TEXT
    cta_tf.paragraphs[0].font.name = BRAND_HEADING_FONT
    cta_tf.paragraphs[0].font.size = Pt(22)
    cta_tf.paragraphs[0].font.bold = True
    cta_tf.paragraphs[0].font.color.rgb = hex_to_rgb("FFFFFF")
    cta_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    cta_tf.word_wrap = False

    # Website URL below button
    url_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.6),
        Inches(6.5), Inches(0.5)
    )
    tf = url_box.text_frame
    tf.anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "base44.com"
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(18)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
    p.alignment = PP_ALIGN.CENTER

    # Save
    output = Path("carousel-cta-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
