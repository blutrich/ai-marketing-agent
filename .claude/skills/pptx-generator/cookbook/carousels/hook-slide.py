#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "carousel-hook-slide"
# format = "carousel"
# dimensions = "square"
# purpose = "Opening slide to grab attention - the scroll-stopper"
# best_for = [
#     "First slide of any carousel",
#     "Bold provocative statements",
#     "Questions that create curiosity",
#     "Controversial or surprising claims",
# ]
# avoid_when = [
#     "Middle of carousel - use single-point-slide instead",
#     "Long explanatory text - save that for body slides",
#     "Multiple points - this is ONE hook only",
# ]
# hook_max_chars = 50
# instructions = [
#     "Hook text should be SHORT and PUNCHY - max 50 characters",
#     "Use confrontational or curiosity-driven language",
#     "Badge text is optional - use for context like 'Just shipped' or 'New'",
#     "Logo appears at bottom for brand recognition",
#     "This slide makes people stop scrolling",
# ]
# ///
"""
LAYOUT: Carousel Hook Slide (Base44 Light Theme)
PURPOSE: Opening slide that stops the scroll - bold statement or question
FORMAT: Square (1:1) for LinkedIn carousels

CUSTOMIZE:
- BRAND_* colors from brand.json
- HOOK_TEXT: The attention-grabbing statement (max 50 chars)
- BADGE_TEXT: Optional context badge (e.g., "Just shipped", "New")
- LOGO_PATH: Path to brand logo
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_logo(slide, logo_path: str, x: float, y: float, width: float = 1.2):
    """Add logo to slide if file exists."""
    logo_file = Path(logo_path)
    if logo_file.exists():
        slide.shapes.add_picture(
            str(logo_file),
            Inches(x), Inches(y),
            width=Inches(width)
        )
        return True
    return False


def main() -> None:
    # === BASE44 BRAND COLORS (Light Theme - Updated 2026-01-24) ===
    BRAND_BG = "E8F4F8"           # Light blue
    BRAND_TEXT = "000000"         # True black text
    BRAND_ACCENT = "FF983B"       # Warm orange
    BRAND_CARD_BG = "FFFFFF"      # White
    BRAND_HEADING_FONT = "Arial"  # Fallback for BC Novatica Cyr

    # === CONTENT ===
    HOOK_TEXT = "Your app now syncs with Google Calendar"
    BADGE_TEXT = "Just shipped"  # Optional - leave empty if not needed
    LOGO_PATH = "brands/base44/assets/logo.png"

    # === SLIDE SETUP (Square 1:1) ===
    prs = Presentation()
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background - light blue
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # White content card - takes up most of slide
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.35), Inches(0.35),
        Inches(6.8), Inches(6.8)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = hex_to_rgb(BRAND_CARD_BG)
    card.line.fill.background()
    card.adjustments[0] = 0.03

    # Orange accent bar at top of card
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.35), Inches(0.35),
        Inches(6.8), Inches(0.1)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    accent_bar.line.fill.background()

    # Badge (optional) - orange pill at top
    if BADGE_TEXT:
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.4), Inches(1.2),
            Inches(2.7), Inches(0.55)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
        badge.line.fill.background()
        badge.adjustments[0] = 0.5  # Fully rounded

        # Badge text - vertically centered
        badge_tf = badge.text_frame
        badge_tf.paragraphs[0].text = BADGE_TEXT
        badge_tf.paragraphs[0].font.name = BRAND_HEADING_FONT
        badge_tf.paragraphs[0].font.size = Pt(18)
        badge_tf.paragraphs[0].font.bold = True
        badge_tf.paragraphs[0].font.color.rgb = hex_to_rgb("FFFFFF")
        badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge_tf.word_wrap = False
        badge.text_frame.anchor = MSO_ANCHOR.MIDDLE

    # Large hook text - centered in card (both H and V)
    hook_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(2.0),
        Inches(6.3), Inches(3.5)
    )
    tf = hook_box.text_frame
    tf.word_wrap = True
    tf.anchor = MSO_ANCHOR.MIDDLE  # Vertical center
    p = tf.paragraphs[0]
    p.text = HOOK_TEXT
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.CENTER  # Horizontal center

    # NOTE: Logo REMOVED from hook slide - logo only on CTA slide
    # This fixes Mac Quick Look rendering bug where text stops rendering
    # when same image appears on multiple slides

    # Save
    output = Path("carousel-hook-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
